import io
import json
import requests
import random
from urllib.parse import quote
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# --- 1. CONFIGURACIÓN DE COLORES INTELIGENTES ---
COLORES_TEMATICOS = {
    "Matemática": RGBColor(25, 55, 109),       
    "Ciencia": RGBColor(40, 80, 40),           
    "Tecnología": RGBColor(70, 70, 70),        
    "Comunicación": RGBColor(160, 20, 20),     
    "Sociales": RGBColor(180, 90, 20),         
    "Arte": RGBColor(140, 20, 140),            
    "Default": RGBColor(0, 51, 102)            
}

# Imágenes de respaldo por si falla la IA (URLs directas de Unsplash)
IMAGENES_BACKUP = [
    "https://images.unsplash.com/photo-1503676260728-1c00da094a0b?q=80&w=800&auto=format&fit=crop", # Educación Gral
    "https://images.unsplash.com/photo-1497633762265-9d179a990aa6?q=80&w=800&auto=format&fit=crop", # Libros
    "https://images.unsplash.com/photo-1509062522246-3755977927d7?q=80&w=800&auto=format&fit=crop", # Salón
    "https://images.unsplash.com/photo-1454165804606-c3d57bc86b40?q=80&w=800&auto=format&fit=crop"  # Trabajo
]

def detectar_color(json_data):
    try:
        subtitulo = json_data.get('slide_1', {}).get('subtitulo', '').lower()
        if "matem" in subtitulo: return COLORES_TEMATICOS["Matemática"]
        if "cien" in subtitulo or "ambien" in subtitulo: return COLORES_TEMATICOS["Ciencia"]
        if "comuni" in subtitulo or "letras" in subtitulo: return COLORES_TEMATICOS["Comunicación"]
        if "social" in subtitulo or "civica" in subtitulo: return COLORES_TEMATICOS["Sociales"]
        if "arte" in subtitulo: return COLORES_TEMATICOS["Arte"]
    except:
        pass
    return COLORES_TEMATICOS["Default"]

# --- 2. DESCARGA DE IMÁGENES (CON RESPALDO) ---
def obtener_imagen_ia(descripcion_en_ingles):
    """
    Intenta IA -> Si falla, usa Backup -> Si falla, devuelve None.
    """
    headers = {'User-Agent': 'Mozilla/5.0'}
    
    # INTENTO 1: IA (Pollinations)
    if descripcion_en_ingles:
        try:
            prompt_seguro = quote(descripcion_en_ingles)
            url = f"https://image.pollinations.ai/prompt/{prompt_seguro}?width=800&height=600&nologo=true&model=flux"
            response = requests.get(url, headers=headers, timeout=15)
            if response.status_code == 200:
                return io.BytesIO(response.content)
        except Exception as e:
            print(f"Fallo IA, usando backup: {e}")

    # INTENTO 2: IMAGEN DE RESPALDO (Backup)
    try:
        url_backup = random.choice(IMAGENES_BACKUP)
        response = requests.get(url_backup, headers=headers, timeout=10)
        if response.status_code == 200:
            return io.BytesIO(response.content)
    except Exception as e:
        print(f"Fallo Backup: {e}")
        
    return None

# --- 3. DIBUJADO DE BARRA ---
def agregar_barra_superior(slide, titulo_texto, color_tema):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color_tema
    shape.line.fill.background()
    
    if slide.shapes.title:
        try: slide.shapes.title.element.getparent().remove(slide.shapes.title.element)
        except: pass
        
    textbox = slide.shapes.add_textbox(Inches(0.2), Inches(0.15), Inches(9.6), Inches(1))
    tf = textbox.text_frame
    tf.vertical_anchor = MSO_SHAPE.RECTANGLE 
    p = tf.paragraphs[0]
    p.text = titulo_texto
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.size = Pt(28) # Letra un pelín más chica para que quepa siempre
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def crear_ppt_desde_data(json_texto):
    prs = Presentation()
    try:
        data = json.loads(json_texto.replace('```json', '').replace('```', '').strip())
    except:
        return None

    color_actual = detectar_color(data)

    # --- SLIDE 1: PORTADA ---
    if 'slide_1' in data:
        info = data['slide_1']
        slide = prs.slides.add_slide(prs.slide_layouts[6]) 
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color_actual
        
        # CORRECCIÓN DE ENCUADRE PORTADA
        # Bajamos el título y le damos más altura
        tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(3))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        p = tf_title.paragraphs[0]
        p.text = info.get('titulo', 'Sin Título')
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(40) 
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER 

        # Subtítulo
        tb_sub = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(1))
        p_sub = tb_sub.text_frame.paragraphs[0]
        p_sub.text = info.get('subtitulo', '')
        p_sub.font.color.rgb = RGBColor(230, 230, 230)
        p_sub.font.size = Pt(24)
        p_sub.alignment = PP_ALIGN.CENTER

    # --- SLIDES 2-7: CONTENIDO ---
    keys = ['slide_2', 'slide_3', 'slide_4', 'slide_5', 'slide_6', 'slide_7']
    
    for key in keys:
        if key in data:
            info = data[key]
            slide = prs.slides.add_slide(prs.slide_layouts[6]) 
            
            # 1. Barra Superior
            agregar_barra_superior(slide, info.get('titulo', ''), color_actual)
            
            # 2. Texto
            tb_body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
            tf = tb_body.text_frame
            tf.word_wrap = True
            
            contenido = info.get('contenido') or info.get('puntos') or ""
            
            def add_p(text, bullet=False):
                p = tf.add_paragraph()
                p.text = str(text)
                p.font.size = Pt(18)
                p.space_after = Pt(10)
                if bullet: p.text = "• " + str(text)

            if isinstance(contenido, list):
                for x in contenido: add_p(x, True)
            else:
                add_p(contenido)

            # 3. IMAGEN (INTENTO ROBUSTO)
            desc = info.get('descripcion_imagen', '')
            if desc:
                img_bytes = obtener_imagen_ia(desc) # Ahora tiene backup automático
                if img_bytes:
                    try:
                        slide.shapes.add_picture(img_bytes, Inches(5.2), Inches(1.5), width=Inches(4.3))
                    except:
                        pass # Si falla backup, se deja en blanco (sin cuadro gris)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
